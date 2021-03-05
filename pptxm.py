from pptx import Presentation
from pptx.util import Inches

def make_ppt(out_ppt_name, content_lst):
    # Presentation()을 일종의 템플릿 객체
    this_prs = Presentation()
    """
    slide_layout[0]는 title, subtitle로 구성된 제목 슬라이드 
    slide_layout[1]는 title, text로 구성된 일반적인 슬라이드 레이아웃
    """
    slide_layout = this_prs.slide_layouts[1] 
    for title, content, img_file_name in content_lst:
        this_slide = this_prs.slides.add_slide(slide_layout)
        shapes = this_slide.shapes
        shapes.title.text = title
        shapes.placeholders[1].text = content
        # placeholders는 개별 slide에 있는 모든 개체를 가져온다고 보면 됨. 
        #shapes.add_picture(img_stream, left, top, height=height)
        #shapes.add_picture(img_file_name, left=Inches(5), top=Inches(10))
        shapes.add_picture(img_file_name, Inches(2.5), Inches(3.2))
        # 변환하지 않고 숫자로 넘기면 잘 되지 않는다. 
    this_prs.save(out_ppt_name)